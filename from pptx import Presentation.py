from pptx import Presentation
from pptx.util import Inches, Pt
import comtypes.client  # for converting PPTX → PDF (Windows only)

# --- Create the PPTX presentation ---
prs = Presentation()
SLIDE_LAYOUT_TITLE = 0
SLIDE_LAYOUT_TITLE_AND_CONTENT = 1

slides_content = [
    ("Cover", "Transforming enterprise procurement with autonomous AI insights"),
    ("Problem", "• High enterprise procurement costs\n• Hidden risks in contracts and vendors\n• Slow manual contract review"),
    ("Solution", "• AI-powered vendor & contract analysis\n• Workflow automation & alerts\n• Actionable recommendations"),
    ("Market Opportunity", "• TAM: $5T global procurement spend\n• SAM: $500B automation-ready spend\n• SOM: $10B high-value enterprise segment"),
    ("Product Architecture", "n8n Orchestration → AI Microservices → Dashboards"),
    ("AI Layer", "• Vendor Knowledge Graph\n• Contract NLP Extraction\n• Pricing Anomaly Detection"),
    ("Automation Layer", "• n8n Workflows\n• Alerts via Slack/Email\n• Executive Reports"),
    ("Security & Compliance", "• RBAC & Tenant Isolation\n• Red-Team Tested\n• SOC2-Ready"),
    ("Technical Stack", "• Python, FastAPI, Gunicorn\n• Docker, Kubernetes, Helm\n• MLflow, Prometheus, Grafana"),
    ("Competitive Advantage", "• Proprietary AI Models\n• End-to-End Automation\n• Enterprise-grade Security"),
    ("Traction / Pilot Results", "• 20 pilot contracts analyzed → $1.2M potential savings\n• Average review time reduced 70%\n• Vendor risk detection accuracy: 95%"),
    ("Go-To-Market Strategy", "• Direct enterprise sales + partnerships\n• Professional services for implementation\n• Pilot → full subscription → upsell analytics"),
    ("Financials Snapshot", "• ARR Target: $10M\n• Cost per Enterprise: $50K/year\n• Projected Customers (Year 5): 100\n• Gross Margin: 75%"),
    ("Team", "• AI Engineer: Platform & ML models\n• DevOps / Kubernetes Specialist\n• Cybersecurity Lead\n• Business Development / GTM"),
    ("Call to Action", "• Investment Ask: $10M\n• Roadmap: Pilot → Deployment → Enterprise Rollout\n• Contact: founder email + @elevated_cuisine_stl")
]

for title, content in slides_content:
    slide = prs.slides.add_slide(SLIDE_LAYOUT_TITLE_AND_CONTENT)
    slide.shapes.title.text = title
    slide.placeholders[1].text_frame.text = content

pptx_filename = "Autonomous_Procurement_Pitch_Deck.pptx"
prs.save(pptx_filename)
print(f"PPTX saved as {pptx_filename}")

# --- Convert PPTX to PDF (Windows only) ---
ppt_app = comtypes.client.CreateObject("PowerPoint.Application")
ppt_app.Visible = 1
presentation = ppt_app.Presentations.Open(pptx_filename)
pdf_filename = "Autonomous_Procurement_Pitch_Deck.pdf"
presentation.SaveAs(pdf_filename, 32)  # 32 = PDF format
presentation.Close()
ppt_app.Quit()
print(f"PDF saved as {pdf_filename}")